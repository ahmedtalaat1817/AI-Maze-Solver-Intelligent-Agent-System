"""Generate PDF and PowerPoint deliverables for the maze project."""

from __future__ import annotations

import csv
from datetime import datetime
from pathlib import Path
import site


ROOT = Path(__file__).resolve().parents[1]
LOCAL_SITE = ROOT / ".venv" / "Lib" / "site-packages"
if LOCAL_SITE.exists():
    site.addsitedir(str(LOCAL_SITE))

from fpdf import FPDF  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


STUDENTS = [
    ("Rany Wael", "230103861"),
    ("Loay Hany", "230105133"),
    ("Magy Romani", "240102851"),
    ("Ereny Magdy", "230104519"),
    ("Ahmed Ehab", "240101817"),
]


def _latest_file(pattern: str):
    matches = sorted(ROOT.glob(pattern), key=lambda path: path.stat().st_mtime)
    return matches[-1] if matches else None


def _latest_results():
    csv_path = _latest_file("outputs/csv/stats_*.csv")
    if not csv_path:
        return None, []

    with csv_path.open(newline="", encoding="utf-8") as handle:
        return csv_path, list(csv.DictReader(handle))


class PDFReport(FPDF):
    def header(self):
        self.set_font("helvetica", "B", 14)
        self.cell(
            0,
            9,
            "Maze Solver Challenge - AI Game Agent",
            0,
            new_x="LMARGIN",
            new_y="NEXT",
            align="C",
        )
        self.set_font("helvetica", "I", 9)
        self.cell(
            0,
            8,
            "El Sewedy University of Technology - CET251 Artificial Intelligence",
            0,
            new_x="LMARGIN",
            new_y="NEXT",
            align="C",
        )
        self.ln(3)

    def footer(self):
        self.set_y(-15)
        self.set_font("helvetica", "I", 8)
        self.cell(0, 10, f"Page {self.page_no()}", 0, new_x="RIGHT", new_y="TOP", align="C")


def generate_pdf():
    output_dir = ROOT / "outputs" / "reports"
    output_dir.mkdir(parents=True, exist_ok=True)
    csv_path, rows = _latest_results()
    chart_path = _latest_file("comparison_chart.png") or _latest_file("outputs/charts/comparison_*.png")

    pdf = PDFReport()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font("helvetica", "B", 22)
    pdf.cell(0, 16, "Final Project Report", 0, new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("helvetica", "", 10)
    pdf.cell(
        0,
        8,
        f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        0,
        new_x="LMARGIN",
        new_y="NEXT",
        align="C",
    )
    pdf.ln(5)

    pdf.set_font("helvetica", "B", 13)
    pdf.cell(0, 8, "Team Members", 0, new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("helvetica", "", 11)
    for name, student_id in STUDENTS:
        pdf.cell(0, 7, f"- {name} ({student_id})", 0, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    sections = [
        (
            "Purpose",
            "This project implements a game-style maze agent that starts from a defined "
            "position, avoids walls, reacts to traps and moving enemies, and reaches a "
            "goal while comparing multiple AI search strategies.",
        ),
        (
            "AI Concepts",
            "The prototype demonstrates agents and environments, uninformed search "
            "(BFS and DFS), informed search (Greedy and A*), cost-aware search (UCS), "
            "partial observability through fog-of-war, and a neural-network risk "
            "predictor for nearby dangerous cells.",
        ),
        (
            "Neural Network Add-on",
            "TrapPredictor trains a small scikit-learn MLP from maze-derived features: "
            "distance to goal, wall density, dead-end status, enemy proximity, known "
            "hazard tiles, and local open space. It outputs a risk probability for "
            "cells on the planned path.",
        ),
        (
            "Evaluation",
            "Each run records success, path steps, weighted path cost, expanded nodes, "
            "runtime, deaths, and average risk. The dashboard exports CSV statistics, "
            "screenshots, and comparison charts for the final demo.",
        ),
    ]

    for title, body in sections:
        pdf.set_font("helvetica", "B", 13)
        pdf.cell(0, 8, title, 0, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("helvetica", "", 11)
        pdf.multi_cell(0, 6, body)
        pdf.ln(3)

    if rows:
        pdf.set_font("helvetica", "B", 13)
        pdf.cell(0, 8, "Latest Experimental Results", 0, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("helvetica", "", 8)
        pdf.cell(24, 6, "Level", border=1)
        pdf.cell(20, 6, "Algo", border=1)
        pdf.cell(18, 6, "OK", border=1)
        pdf.cell(18, 6, "Steps", border=1)
        pdf.cell(18, 6, "Cost", border=1)
        pdf.cell(22, 6, "Nodes", border=1)
        pdf.cell(24, 6, "Time(s)", border=1)
        pdf.cell(22, 6, "Risk", border=1, new_x="LMARGIN", new_y="NEXT")

        for row in rows[-10:]:
            pdf.cell(24, 6, row.get("level", "")[:12], border=1)
            pdf.cell(20, 6, row.get("algo", ""), border=1)
            pdf.cell(18, 6, str(row.get("success", "")), border=1)
            pdf.cell(18, 6, row.get("path_len", ""), border=1)
            pdf.cell(18, 6, row.get("path_cost", ""), border=1)
            pdf.cell(22, 6, row.get("nodes", ""), border=1)
            pdf.cell(24, 6, str(row.get("time", ""))[:7], border=1)
            risk = row.get("avg_risk", "")
            risk_text = f"{float(risk):.0%}" if risk else ""
            pdf.cell(22, 6, risk_text, border=1, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(3)
        pdf.set_font("helvetica", "I", 8)
        pdf.cell(0, 6, f"Source CSV: {csv_path}", 0, new_x="LMARGIN", new_y="NEXT")

    if chart_path and chart_path.exists():
        pdf.add_page()
        pdf.set_font("helvetica", "B", 13)
        pdf.cell(0, 8, "Comparison Chart", 0, new_x="LMARGIN", new_y="NEXT")
        pdf.image(str(chart_path), x=12, w=186)

    output_path = output_dir / "Maze_Solver_Report.pdf"
    pdf.output(output_path)
    print(f"[PDF] Generated: {output_path}")


def _add_title(slide, title_text, subtitle_text=None):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 120, 130)
    if subtitle_text and len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text


def generate_pptx():
    output_dir = ROOT / "outputs" / "reports"
    output_dir.mkdir(parents=True, exist_ok=True)
    _, rows = _latest_results()
    chart_path = _latest_file("comparison_chart.png") or _latest_file("outputs/charts/comparison_*.png")

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(
        slide,
        "Maze Solver Challenge",
        "AI Game Agent | BFS, DFS, UCS, Greedy, A*, and MLP Risk Prediction",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_title(slide, "Project Scope")
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Working Python prototype with a Pygame dashboard"
    for text in [
        "At least three mazes plus procedural level generation",
        "Five search algorithms with runtime and node-expansion metrics",
        "Moving enemies, weighted terrain, fog-of-war, and trap risk prediction",
        "CSV, screenshot, PDF, PPTX, and chart outputs",
    ]:
        paragraph = body.add_paragraph()
        paragraph.text = text
        paragraph.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_title(slide, "Algorithm Comparison")
    body = slide.shapes.placeholders[1].text_frame
    body.text = "BFS and DFS provide baseline uninformed search behavior"
    for text in [
        "UCS minimizes weighted path cost across sand, water, and hazard cells",
        "Greedy search uses only the goal-distance heuristic",
        "A* combines path cost and heuristic distance for efficient optimal search",
    ]:
        paragraph = body.add_paragraph()
        paragraph.text = text
        paragraph.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_title(slide, "Neural Network Risk Predictor")
    body = slide.shapes.placeholders[1].text_frame
    body.text = "TrapPredictor trains a small MLP on maze-derived synthetic labels"
    for text in [
        "Inputs: goal distance, wall density, dead ends, enemy proximity, hazards, open space",
        "Output: danger probability for each cell in the candidate path",
        "GUI visualizes safe cells in green and risky cells in red",
    ]:
        paragraph = body.add_paragraph()
        paragraph.text = text
        paragraph.level = 1

    if chart_path and chart_path.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _add_title(slide, "Latest Results Chart")
        slide.shapes.add_picture(str(chart_path), Inches(0.45), Inches(1.2), width=Inches(9.1))

    if rows:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        _add_title(slide, "Latest Run Summary")
        body = slide.shapes.placeholders[1].text_frame
        body.text = "Most recent recorded results:"
        for row in rows[-5:]:
            paragraph = body.add_paragraph()
            paragraph.text = (
                f"{row.get('level')} - {row.get('algo')}: "
                f"steps {row.get('path_len')}, cost {row.get('path_cost')}, "
                f"nodes {row.get('nodes')}, success {row.get('success')}"
            )
            paragraph.level = 1

    output_path = output_dir / "Maze_Solver_Presentation.pptx"
    prs.save(output_path)
    print(f"[PPTX] Generated: {output_path}")


if __name__ == "__main__":
    generate_pdf()
    generate_pptx()
